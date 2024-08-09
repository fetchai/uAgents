
import json
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import random
import tkinter as tk
from tkinter import simpledialog, messagebox
from open_ai_assistant import OpenAIProcessor
from uagents import Agent, Context

# Initialize OpenAI processor
openai_processor = OpenAIProcessor(api_key="your_api_key_opneai")

# Create the agent
ppt_agent = Agent(name="ppt_agent", seed="ppt_agent_seed")

@ppt_agent.on_event("startup")
async def startup_message(ctx: Context):
    ctx.logger.info(f"Agent {ppt_agent.name} has started with address {ppt_agent.address}.")
    input_data = get_user_input_via_gui()
    ctx.logger.info("got the input data")
    create_ppt_from_input(input_data["topic"], input_data["description"], input_data["slides_number"])
    ctx.logger.info("Presentation saved as 'output_presentation.pptx'")

# @ppt_agent.on_interval(period=5.0)
# async def log_agent_activity(ctx: Context):
#     ctx.logger.info("Agent is active and waiting for user inputs.")

def download_random_image():

    random_width = random.randint(1920, 2560)
    random_height = random.randint(1080, 1440)
    image_url = f"https://picsum.photos/{random_width}/{random_height}"
    response = requests.get(image_url)
    return BytesIO(response.content)

def add_background_image(slide, image_stream):
    slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

def create_text_box_with_blended_background(slide, left, top, width, height, text, font_size):
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = background_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    fill.fore_color.brightness = 0.5

    background_shape.line.color.rgb = RGBColor(255, 255, 255)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.alignment = PP_ALIGN.LEFT

def create_ppt_from_input(topic, description, slides_number):
    slides_content = openai_processor.get_assistant_response(topic, description, slides_number)
    slides_content = json.loads(slides_content)

    global prs
    prs = Presentation()

    for slide_num, slide_data in slides_content.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        image_stream = download_random_image()
        add_background_image(slide, image_stream)
        
        main_topic = slide_data.get("main_topic", "Main Topic").strip()
        create_text_box_with_blended_background(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(1.5), main_topic, font_size=28)

        top = Inches(2)
        left_margin = Inches(0.75)
        right_margin = prs.slide_width - Inches(1.5)

        subtopics = slide_data.get("subtopics", [])
        for subtopic in subtopics:
            subtopic_title = subtopic.get("title", "Subtopic").strip()
            content = subtopic.get("content", "Content goes here.").strip()
            
            create_text_box_with_blended_background(slide, left_margin, top, right_margin - left_margin, Inches(0.5), subtopic_title, font_size=24)
            top += Inches(0.6)
            
            create_text_box_with_blended_background(slide, left_margin + Inches(0.25), top, right_margin - left_margin - Inches(0.5), Inches(1), content, font_size=20)
            top += Inches(1.2)

    prs.save('output_presentation.pptx')


def get_user_input_via_gui():
    root = tk.Tk()
    root.withdraw()  # Hide the root window

    topic = simpledialog.askstring("Input", "Enter the topic of the presentation:", parent=root)
    description = simpledialog.askstring("Input", "Enter a brief description of the presentation:", parent=root)
    slides_number = simpledialog.askinteger("Input", "Enter the number of slides you want to create:", parent=root)

    if not topic or not description or slides_number is None:
        messagebox.showerror("Error", "All inputs must be provided!")
        root.quit()
        return None

    return {"topic": topic, "description": description, "slides_number": slides_number}

if __name__ == "__main__":

    
    ppt_agent.run()
